#!/usr/bin/env python3
"""Generate Industry AI Flow Capstone Presentation (.pptx)

Team: Jason Niu, Jack Si, Angel Daniel Bustamante Perez
Program: Integrated Artificial Intelligence, SAIT
Date: March 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette (Professional Light Mode) ─────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE = RGBColor(0xF8, 0xF9, 0xFA)
LIGHT_GRAY = RGBColor(0xE9, 0xEC, 0xEF)
MID_GRAY = RGBColor(0x6C, 0x75, 0x7D)
DARK_TEXT = RGBColor(0x21, 0x25, 0x29)
ACCENT_BLUE = RGBColor(0x00, 0x6D, 0xC6)  # Primary accent
ACCENT_TEAL = RGBColor(0x00, 0x96, 0x88)  # Secondary accent
ACCENT_ORANGE = RGBColor(0xF5, 0x7C, 0x00)  # Highlight
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
LIGHT_BLUE_BG = RGBColor(0xE3, 0xF2, 0xFD)
LIGHT_TEAL_BG = RGBColor(0xE0, 0xF2, 0xF1)
LIGHT_ORANGE_BG = RGBColor(0xFF, 0xF3, 0xE0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helper Functions ─────────────────────────────────────────────────
def add_bg(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.2):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(font_size * 0.3)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf


def add_multi_text(slide, left, top, width, height, lines, default_size=16,
                   font_name="Calibri", line_spacing=1.3):
    """Add textbox with multiple styled lines. Each line is (text, size, bold, color)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        bold = line_data[2] if len(line_data) > 2 else False
        color = line_data[3] if len(line_data) > 3 else DARK_TEXT
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(size * 0.4)
        p.line_spacing = Pt(size * line_spacing)
    return tf


def add_rect(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_lines,
             accent_color=ACCENT_BLUE, bg_color=LIGHT_BLUE_BG):
    """Add a styled card with accent top border."""
    # Card background
    card = add_rect(slide, left, top, width, height, bg_color)
    # Accent top bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # Title
    add_textbox(slide, left + Inches(0.25), top + Inches(0.15),
                width - Inches(0.5), Inches(0.4),
                title, font_size=16, bold=True, color=accent_color)
    # Body
    y_offset = Inches(0.6)
    for line in body_lines:
        add_textbox(slide, left + Inches(0.25), top + y_offset,
                    width - Inches(0.5), Inches(0.35),
                    line, font_size=13, color=DARK_TEXT)
        y_offset += Inches(0.3)


def add_slide_number(slide, num, total=14):
    """Add slide number at bottom right."""
    add_textbox(slide, Inches(12.0), Inches(7.05), Inches(1.1), Inches(0.35),
                f"{num} / {total}", font_size=10, color=MID_GRAY,
                alignment=PP_ALIGN.RIGHT)


def add_speaker_tag(slide, name):
    """Add speaker indicator at bottom left."""
    tag = add_rect(slide, Inches(0.3), Inches(7.0), Inches(1.8), Inches(0.35),
                   LIGHT_GRAY, border_color=MID_GRAY)
    add_textbox(slide, Inches(0.35), Inches(7.02), Inches(1.7), Inches(0.3),
                f"Speaker: {name}", font_size=9, color=MID_GRAY,
                alignment=PP_ALIGN.CENTER)


def add_section_header_line(slide, y=Inches(1.45)):
    """Add a thin accent line under slide title."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(2.0), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ══════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide1, WHITE)

# Top accent bar
bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

# Main title
add_textbox(slide1, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
            "Industry AI Flow", font_size=52, bold=True, color=DARK_TEXT,
            alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide1, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.6),
            "AI-Powered Intelligence for the Construction Industry",
            font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Divider line
div = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(3.8), Inches(3.3), Inches(0.03)
)
div.fill.solid()
div.fill.fore_color.rgb = LIGHT_GRAY
div.line.fill.background()

# Team & Program info
add_multi_text(slide1, Inches(1.5), Inches(4.2), Inches(10.3), Inches(2.5), [
    ("SAIT  |  Integrated Artificial Intelligence  |  Capstone Project", 14, False, MID_GRAY, PP_ALIGN.CENTER),
    ("", 8),
    ("Angel Daniel Bustamante Perez    |    Jason Niu    |    Jack Si", 16, False, DARK_TEXT, PP_ALIGN.CENTER),
    ("", 8),
    ("Instructor: Reeta", 13, False, MID_GRAY, PP_ALIGN.CENTER),
    ("March 2026", 13, False, MID_GRAY, PP_ALIGN.CENTER),
])

add_slide_number(slide1, 1)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2: Agenda / Outline  (Angel)
# ══════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)

add_textbox(slide2, Inches(0.6), Inches(0.5), Inches(6), Inches(0.7),
            "Agenda", font_size=36, bold=True, color=DARK_TEXT)
add_section_header_line(slide2)

agenda_items = [
    ("1", "Problem & Motivation", "Why AI for construction?", ACCENT_ORANGE),
    ("2", "Our Solution", "Three AI-powered capabilities", ACCENT_BLUE),
    ("3", "System Architecture", "How it all fits together", ACCENT_TEAL),
    ("4", "AI/ML Deep Dive", "RAG, cost prediction, code analysis", ACCENT_PURPLE),
    ("5", "Results & Evaluation", "Metrics and achievements", ACCENT_GREEN),
    ("6", "Conclusion & Future Work", "Key takeaways", MID_GRAY),
]

for i, (num, title, desc, color) in enumerate(agenda_items):
    y = Inches(2.0) + i * Inches(0.85)
    # Number circle
    circle = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_textbox(slide2, Inches(1.5), y + Inches(0.05), Inches(0.5), Inches(0.4),
                num, font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title + Description
    add_textbox(slide2, Inches(2.3), y + Inches(0.0), Inches(4), Inches(0.3),
                title, font_size=18, bold=True, color=DARK_TEXT)
    add_textbox(slide2, Inches(2.3), y + Inches(0.3), Inches(4), Inches(0.25),
                desc, font_size=12, color=MID_GRAY)

# Right side: speaker assignment box
add_rect(slide2, Inches(7.8), Inches(2.0), Inches(4.5), Inches(4.5), NEAR_WHITE, border_color=LIGHT_GRAY)
add_textbox(slide2, Inches(8.0), Inches(2.15), Inches(4.0), Inches(0.4),
            "Speaker Assignments", font_size=16, bold=True, color=ACCENT_BLUE)

speakers = [
    ("Angel", "Introduction, Process, Conclusion", ACCENT_ORANGE),
    ("Jack", "Industry Problem, Domain Data", ACCENT_TEAL),
    ("Jason", "Architecture, AI/ML, Demo", ACCENT_BLUE),
]
for i, (name, role, color) in enumerate(speakers):
    y = Inches(2.8) + i * Inches(1.1)
    add_rect(slide2, Inches(8.1), y, Inches(4.0), Inches(0.85), WHITE, border_color=color)
    add_textbox(slide2, Inches(8.3), y + Inches(0.08), Inches(3.5), Inches(0.35),
                name, font_size=16, bold=True, color=color)
    add_textbox(slide2, Inches(8.3), y + Inches(0.45), Inches(3.5), Inches(0.3),
                role, font_size=12, color=MID_GRAY)

add_speaker_tag(slide2, "Angel")
add_slide_number(slide2, 2)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3: Problem & Motivation  (Jack)
# ══════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)

add_textbox(slide3, Inches(0.6), Inches(0.5), Inches(8), Inches(0.7),
            "The Problem: Construction Meets AI", font_size=36, bold=True, color=DARK_TEXT)
add_section_header_line(slide3)

# Left: Pain points
pain_points = [
    ("Scattered Knowledge", "Safety codes, regulations, and standards spread across hundreds of PDF documents — hard to search, easy to miss"),
    ("Cost Overruns", "Large construction projects frequently exceed budget (McKinsey, 2017). Manual estimates are slow and inconsistent"),
    ("Data Underutilization", "Project data sits in spreadsheets. Teams lack tools to extract insights without data science skills"),
]

for i, (title, desc) in enumerate(pain_points):
    y = Inches(1.9) + i * Inches(1.6)
    # Warning icon (orange dot)
    dot = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_ORANGE
    dot.line.fill.background()
    add_textbox(slide3, Inches(0.8), y + Inches(0.02), Inches(0.3), Inches(0.3),
                "!", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide3, Inches(1.3), y, Inches(5.2), Inches(0.35),
                title, font_size=18, bold=True, color=DARK_TEXT)
    add_textbox(slide3, Inches(1.3), y + Inches(0.4), Inches(5.2), Inches(0.9),
                desc, font_size=13, color=MID_GRAY)

# Right: Opportunity box
add_rect(slide3, Inches(7.2), Inches(1.9), Inches(5.3), Inches(4.8), LIGHT_BLUE_BG)
add_textbox(slide3, Inches(7.5), Inches(2.1), Inches(4.7), Inches(0.4),
            "The Opportunity", font_size=20, bold=True, color=ACCENT_BLUE)

opportunities = [
    "AI can instantly search and cite relevant safety codes from a vectorized document library",
    "ML models can predict cost overruns early, before they become project failures",
    "Automated code generation lets non-technical users analyze datasets with natural language",
    "One unified platform replaces fragmented tools",
]
for i, opp in enumerate(opportunities):
    y = Inches(2.7) + i * Inches(1.0)
    # Checkmark
    check = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), y, Inches(0.25), Inches(0.25))
    check.fill.solid()
    check.fill.fore_color.rgb = ACCENT_GREEN
    check.line.fill.background()
    add_textbox(slide3, Inches(7.5), y - Inches(0.03), Inches(0.25), Inches(0.25),
                "+", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide3, Inches(7.9), y - Inches(0.05), Inches(4.3), Inches(0.85),
                opp, font_size=13, color=DARK_TEXT)

add_speaker_tag(slide3, "Jack")
add_slide_number(slide3, 3)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4: Domain Data  (Jack)
# ══════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)

add_textbox(slide4, Inches(0.6), Inches(0.5), Inches(8), Inches(0.7),
            "Construction Domain Data", font_size=36, bold=True, color=DARK_TEXT)
add_section_header_line(slide4)

# Left: Documents
add_card(slide4, Inches(0.6), Inches(1.9), Inches(5.8), Inches(2.5),
         "RAG Knowledge Base  (9+ Indexed Documents)",
         [
             "Ontario Building Code Act 1992",
             "OSHA 29 CFR 1926 Construction Safety",
             "Ontario Reg 213/91 Construction Projects",
             "BuildingSmart IFC 4.3 Schema Specs",
             "Industry safety guides, procedures, manuals",
         ],
         accent_color=ACCENT_BLUE, bg_color=LIGHT_BLUE_BG)

# Left: Dataset
add_card(slide4, Inches(0.6), Inches(4.7), Inches(5.8), Inches(2.2),
         "Cost Estimation Dataset  (10,000 Records)",
         [
             "15 numeric + 2 categorical features (42 after encoding)",
             "Features: project_type, sqft, floors, location, risk_score...",
             "Curated by Jack (construction domain expert)",
             "Covers 12 project types, 15 Canadian locations",
         ],
         accent_color=ACCENT_TEAL, bg_color=LIGHT_TEAL_BG)

# Right: Data pipeline
add_rect(slide4, Inches(6.8), Inches(1.9), Inches(5.8), Inches(5.0), NEAR_WHITE, border_color=LIGHT_GRAY)
add_textbox(slide4, Inches(7.0), Inches(2.05), Inches(5.4), Inches(0.4),
            "ETL Pipeline", font_size=18, bold=True, color=ACCENT_BLUE)

etl_steps = [
    ("EXTRACT", "Upload PDF/images/CSV via web UI\nPaddleOCR for scanned documents", ACCENT_ORANGE),
    ("TRANSFORM", "Text -> 512-char chunks (128 overlap)\nEmbed with nomic-embed-text-v1.5 (768-dim)\nBM25 keyword index built in-memory", ACCENT_BLUE),
    ("LOAD", "Chunks + vectors -> PostgreSQL + pgvector\nIVFFlat index for fast similarity search\nCost model -> JSON artifact on filesystem", ACCENT_GREEN),
]

for i, (phase, desc, color) in enumerate(etl_steps):
    y = Inches(2.65) + i * Inches(1.4)
    # Phase label
    label_bg = add_rect(slide4, Inches(7.1), y, Inches(1.5), Inches(0.35), color)
    add_textbox(slide4, Inches(7.1), y + Inches(0.02), Inches(1.5), Inches(0.3),
                phase, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Description
    for j, line in enumerate(desc.split("\n")):
        add_textbox(slide4, Inches(8.8), y + j * Inches(0.3), Inches(3.5), Inches(0.3),
                    line, font_size=11, color=DARK_TEXT)
    # Arrow down (except last)
    if i < 2:
        arrow_y = y + Inches(1.15)
        add_textbox(slide4, Inches(7.7), arrow_y, Inches(0.5), Inches(0.3),
                    "v", font_size=14, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_speaker_tag(slide4, "Jack")
add_slide_number(slide4, 4)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5: Solution Overview — 3 Capabilities  (Jason)
# ══════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, WHITE)

add_textbox(slide5, Inches(0.6), Inches(0.5), Inches(10), Inches(0.7),
            "Our Solution: Three AI-Powered Capabilities", font_size=36, bold=True, color=DARK_TEXT)
add_section_header_line(slide5)

capabilities = [
    ("RAG Knowledge QA", "Ask questions about construction\ndocuments, get cited answers",
     ["Hybrid search: BM25 + Vector + RRF", "BGE cross-encoder reranking",
      "Source citations on every answer", "Suggested follow-up questions"],
     ACCENT_BLUE, LIGHT_BLUE_BG),

    ("Cost Estimation", "Predict construction project\ncost overruns with ML",
     ["Ridge regression (R\u00b2 = 0.989)", "10,000-record training dataset",
      "12 project types, 15 locations", "Reasonableness validation"],
     ACCENT_TEAL, LIGHT_TEAL_BG),

    ("Dynamic Data Analysis", "Upload datasets, ask questions,\nget charts & insights",
     ["Cloud LLM code generation", "Docker sandbox execution",
      "Dual fallback: Zhipu -> Gemini", "Security-hardened (36 rounds)"],
     ACCENT_PURPLE, LIGHT_ORANGE_BG),
]

for i, (title, subtitle, features, color, bg) in enumerate(capabilities):
    x = Inches(0.6) + i * Inches(4.2)
    # Card
    add_rect(slide5, x, Inches(1.9), Inches(3.8), Inches(5.0), bg, border_color=color)
    # Number badge
    badge = slide5.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), Inches(2.1), Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    add_textbox(slide5, x + Inches(0.2), Inches(2.13), Inches(0.5), Inches(0.45),
                str(i + 1), font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide5, x + Inches(0.8), Inches(2.1), Inches(2.8), Inches(0.4),
                title, font_size=18, bold=True, color=color)
    # Subtitle
    for j, line in enumerate(subtitle.split("\n")):
        add_textbox(slide5, x + Inches(0.3), Inches(2.65) + j * Inches(0.28), Inches(3.2), Inches(0.3),
                    line, font_size=12, color=MID_GRAY)
    # Divider
    div = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x + Inches(0.3), Inches(3.35), Inches(3.2), Inches(0.02)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = LIGHT_GRAY
    div.line.fill.background()
    # Features
    for j, feat in enumerate(features):
        y = Inches(3.6) + j * Inches(0.65)
        # Bullet dot
        dot = slide5.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.35), y + Inches(0.05),
                                       Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_textbox(slide5, x + Inches(0.6), y - Inches(0.03), Inches(3.0), Inches(0.55),
                    feat, font_size=13, color=DARK_TEXT)

add_speaker_tag(slide5, "Jason")
add_slide_number(slide5, 5)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6: System Architecture Overview  (Jason)
# ══════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, WHITE)

add_textbox(slide6, Inches(0.6), Inches(0.5), Inches(10), Inches(0.7),
            "System Architecture", font_size=36, bold=True, color=DARK_TEXT)
add_section_header_line(slide6)

# 6-layer architecture
layers = [
    ("UI Layer", "Next.js + React", ACCENT_BLUE, LIGHT_BLUE_BG),
    ("API Gateway", "FastAPI + Auth + Rate Limiting", ACCENT_TEAL, LIGHT_TEAL_BG),
    ("Business Services", "Intent Classifier + Workflow Orchestrator + Router", ACCENT_ORANGE, LIGHT_ORANGE_BG),
    ("AI Runtime", "RAG Engine  |  Cost Estimator  |  Code Sandbox", ACCENT_PURPLE, RGBColor(0xF3, 0xE5, 0xF5)),
    ("Data Storage", "PostgreSQL + pgvector  |  BM25 Index  |  Model Artifacts", ACCENT_GREEN, RGBColor(0xE8, 0xF5, 0xE9)),
    ("Security & Observability", "Prometheus  |  Audit Logging  |  Tenant Isolation", MID_GRAY, NEAR_WHITE),
]

for i, (name, tech, color, bg) in enumerate(layers):
    y = Inches(1.8) + i * Inches(0.88)
    # Layer bar
    add_rect(slide6, Inches(0.8), y, Inches(7.2), Inches(0.75), bg, border_color=color)
    # Layer name
    add_textbox(slide6, Inches(1.0), y + Inches(0.05), Inches(2.5), Inches(0.35),
                name, font_size=15, bold=True, color=color)
    # Tech detail
    add_textbox(slide6, Inches(3.2), y + Inches(0.05), Inches(4.5), Inches(0.6),
                tech, font_size=12, color=DARK_TEXT)
    # Arrow down (between layers)
    if i < 5:
        arrow_y = y + Inches(0.72)
        add_textbox(slide6, Inches(4.2), arrow_y, Inches(0.5), Inches(0.2),
                    "|", font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Right: Tech stack box
add_rect(slide6, Inches(8.5), Inches(1.8), Inches(4.3), Inches(5.3), NEAR_WHITE, border_color=LIGHT_GRAY)
add_textbox(slide6, Inches(8.7), Inches(1.95), Inches(3.8), Inches(0.35),
            "Tech Stack", font_size=18, bold=True, color=ACCENT_BLUE)

stack_items = [
    ("LLM", "Qwen 3.5:4b via Ollama (Metal GPU)"),
    ("Embeddings", "nomic-embed-text-v1.5 (768-dim)"),
    ("Vector DB", "PostgreSQL 16 + pgvector (IVFFlat)"),
    ("Reranker", "BGE-reranker-base cross-encoder"),
    ("OCR", "PaddleOCR (Python 3.13)"),
    ("Backend", "FastAPI + LangChain 1.0"),
    ("Frontend", "Next.js + TypeScript"),
    ("Sandbox", "Docker (security-hardened)"),
    ("Cloud LLM", "Zhipu GLM-4 / Google Gemini"),
]

for i, (label, value) in enumerate(stack_items):
    y = Inches(2.55) + i * Inches(0.48)
    add_textbox(slide6, Inches(8.7), y, Inches(1.4), Inches(0.3),
                label, font_size=11, bold=True, color=ACCENT_BLUE)
    add_textbox(slide6, Inches(10.1), y, Inches(2.5), Inches(0.4),
                value, font_size=11, color=DARK_TEXT)

add_speaker_tag(slide6, "Jason")
add_slide_number(slide6, 6)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7: Intent Classification Pipeline  (Jason)
# ══════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, WHITE)

add_textbox(slide7, Inches(0.6), Inches(0.5), Inches(10), Inches(0.7),
            "AI Innovation: Intent Classification Pipeline", font_size=36, bold=True, color=DARK_TEXT)
add_section_header_line(slide7)

# Left: 11-node StateGraph
add_textbox(slide7, Inches(0.6), Inches(1.7), Inches(6), Inches(0.4),
            "11-Node LangGraph StateGraph", font_size=18, bold=True, color=ACCENT_BLUE)

add_textbox(slide7, Inches(0.6), Inches(2.15), Inches(6), Inches(0.6),
            "A multi-step graph that classifies user intent, handles ambiguous queries\nwith clarification loops, and routes to the correct AI capability.",
            font_size=13, color=MID_GRAY)

# Flow diagram (simplified as boxes + arrows)
flow_nodes = [
    ("User Input", ACCENT_BLUE),
    ("Intent Classify", ACCENT_PURPLE),
    ("Confidence Check", ACCENT_ORANGE),
]
for i, (label, color) in enumerate(flow_nodes):
    x = Inches(0.8) + i * Inches(2.2)
    add_rect(slide7, x, Inches(3.0), Inches(1.8), Inches(0.6), color)
    add_textbox(slide7, x, Inches(3.08), Inches(1.8), Inches(0.5),
                label, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < 2:
        add_textbox(slide7, x + Inches(1.8), Inches(3.08), Inches(0.4), Inches(0.5),
                    "->", font_size=14, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Branch arrows
add_textbox(slide7, Inches(5.2), Inches(3.7), Inches(2), Inches(0.35),
            ">= 0.8 confidence", font_size=11, color=ACCENT_GREEN)
add_textbox(slide7, Inches(5.2), Inches(4.1), Inches(2), Inches(0.35),
            "< 0.8 -> Clarify (max 2 rounds)", font_size=11, color=ACCENT_ORANGE)

# Three routing targets
routes = [
    ("RAG QA", ACCENT_BLUE),
    ("Cost Estimation", ACCENT_TEAL),
    ("Data Analysis", ACCENT_PURPLE),
]
for i, (label, color) in enumerate(routes):
    x = Inches(0.8) + i * Inches(2.2)
    add_rect(slide7, x, Inches(4.8), Inches(1.8), Inches(0.5), color)
    add_textbox(slide7, x, Inches(4.85), Inches(1.8), Inches(0.4),
                label, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide7, Inches(0.8), Inches(5.5), Inches(5.5), Inches(0.35),
            "Routes to correct capability with reformulated query + extracted keywords",
            font_size=12, color=MID_GRAY)

# Right: 10-node execution pipeline
add_rect(slide7, Inches(7.5), Inches(1.7), Inches(5.3), Inches(5.2), NEAR_WHITE, border_color=LIGHT_GRAY)
add_textbox(slide7, Inches(7.7), Inches(1.85), Inches(4.8), Inches(0.35),
            "10-Node Execution Pipeline", font_size=18, bold=True, color=ACCENT_TEAL)

pipeline_nodes = [
    "1. Intent", "2. Safety", "3. Cost Est.", "4. Retrieval", "5. Rerank",
    "6. Prompt", "7. Route", "8. Code Exec", "9. Response", "10. Groundedness"
]
for i, node in enumerate(pipeline_nodes):
    row = i // 2
    col = i % 2
    x = Inches(7.8) + col * Inches(2.5)
    y = Inches(2.4) + row * Inches(0.85)
    node_color = ACCENT_TEAL if i % 2 == 0 else ACCENT_BLUE
    add_rect(slide7, x, y, Inches(2.2), Inches(0.55), node_color)
    add_textbox(slide7, x, y + Inches(0.08), Inches(2.2), Inches(0.4),
                node, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide7, Inches(7.7), Inches(6.55), Inches(4.8), Inches(0.3),
            "Fixed-order pipeline ensures consistent, auditable processing",
            font_size=11, color=MID_GRAY)

add_speaker_tag(slide7, "Jason")
add_slide_number(slide7, 7)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8: RAG Deep Dive  (Jason)
# ══════════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, WHITE)

add_textbox(slide8, Inches(0.6), Inches(0.5), Inches(10), Inches(0.7),
            "RAG Pipeline: Hybrid Retrieval", font_size=36, bold=True, color=DARK_TEXT)
add_section_header_line(slide8)

# Pipeline flow (horizontal)
rag_steps = [
    ("Query", "User question\n+ keywords", ACCENT_BLUE),
    ("BM25", "Keyword\nmatching", ACCENT_ORANGE),
    ("Vector", "Semantic\nsimilarity", ACCENT_PURPLE),
    ("RRF", "Reciprocal\nRank Fusion", ACCENT_TEAL),
    ("Rerank", "BGE cross-\nencoder", ACCENT_GREEN),
    ("LLM", "Qwen 3.5\ngenerate", ACCENT_BLUE),
    ("Answer", "Cited\nresponse", MID_GRAY),
]

for i, (label, desc, color) in enumerate(rag_steps):
    x = Inches(0.3) + i * Inches(1.82)
    y = Inches(2.0)
    # Box
    add_rect(slide8, x, y, Inches(1.5), Inches(1.5), color)
    add_textbox(slide8, x, y + Inches(0.15), Inches(1.5), Inches(0.4),
                label, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    for j, line in enumerate(desc.split("\n")):
        add_textbox(slide8, x, y + Inches(0.55) + j * Inches(0.3), Inches(1.5), Inches(0.3),
                    line, font_size=11, color=RGBColor(0xE0, 0xE0, 0xE0), alignment=PP_ALIGN.CENTER)
    # Arrow
    if i < 6:
        add_textbox(slide8, x + Inches(1.5), y + Inches(0.45), Inches(0.35), Inches(0.4),
                    ">", font_size=18, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom: Key design decisions
add_textbox(slide8, Inches(0.6), Inches(3.9), Inches(5), Inches(0.4),
            "Key Design Decisions", font_size=18, bold=True, color=ACCENT_BLUE)

decisions = [
    ("Hybrid Search (BM25 + Vector)", "Exact keyword matches + semantic understanding = better recall than either alone"),
    ("Cross-Encoder Reranking", "BGE-reranker-base re-scores top candidates for precision (not just cosine similarity)"),
    ("512-char Chunks, 128 Overlap", "Balances context completeness vs. embedding quality; overlap prevents boundary information loss"),
    ("Source Citations", "Every answer includes document name + chunk reference for traceability and trust"),
]

for i, (title, desc) in enumerate(decisions):
    y = Inches(4.4) + i * Inches(0.72)
    add_textbox(slide8, Inches(0.8), y, Inches(3.5), Inches(0.3),
                title, font_size=13, bold=True, color=DARK_TEXT)
    add_textbox(slide8, Inches(4.3), y, Inches(8.5), Inches(0.6),
                desc, font_size=12, color=MID_GRAY)

add_speaker_tag(slide8, "Jason")
add_slide_number(slide8, 8)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9: Cost Estimation ML  (Jason)
# ══════════════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, WHITE)

add_textbox(slide9, Inches(0.6), Inches(0.5), Inches(10), Inches(0.7),
            "Cost Estimation: Ridge Regression Model", font_size=36, bold=True, color=DARK_TEXT)
add_section_header_line(slide9)

# Left: Model details
add_card(slide9, Inches(0.6), Inches(1.9), Inches(5.5), Inches(2.5),
         "Model Configuration",
         [
             "Algorithm: Ridge Regression (L2 regularization)",
             "Training: 5-fold cross-validation on 10,000 records",
             "Features: 15 numeric + 2 categorical (42 after one-hot)",
             "Normalization: Z-score standardization",
             "Output: Predicted actual cost + overrun percentage",
         ],
         accent_color=ACCENT_TEAL, bg_color=LIGHT_TEAL_BG)

# Left bottom: Top drivers
add_card(slide9, Inches(0.6), Inches(4.7), Inches(5.5), Inches(2.2),
         "Key Cost Features (Model Inputs)",
         [
             "Change Orders — strong cost overrun indicator",
             "Risk Score — project complexity measure",
             "Subcontractor Count — coordination factor",
             "Contractor Rating — quality/reliability signal",
         ],
         accent_color=ACCENT_ORANGE, bg_color=LIGHT_ORANGE_BG)

# Right: Metrics comparison
add_rect(slide9, Inches(6.5), Inches(1.9), Inches(6.2), Inches(5.0), NEAR_WHITE, border_color=LIGHT_GRAY)
add_textbox(slide9, Inches(6.7), Inches(2.05), Inches(5.8), Inches(0.4),
            "Model vs. Baseline Performance", font_size=18, bold=True, color=ACCENT_BLUE)

# Table header
metrics_header = ["Metric", "Baseline", "Our Model", "Improvement"]
for j, h in enumerate(metrics_header):
    x = Inches(6.8) + j * Inches(1.45)
    add_rect(slide9, x, Inches(2.6), Inches(1.35), Inches(0.4), ACCENT_BLUE)
    add_textbox(slide9, x, Inches(2.63), Inches(1.35), Inches(0.35),
                h, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Table rows
metrics_data = [
    ("MAE", "$4.06M", "$2.58M", "36% better"),
    ("MAPE", "8.59%", "6.18%", "28% better"),
    ("R\u00b2", "0.975", "0.989", "+1.4%"),
    ("RMSE", "N/A", "$9.83M", "--"),
]

for i, (metric, baseline, model, imp) in enumerate(metrics_data):
    y = Inches(3.05) + i * Inches(0.55)
    row_bg = WHITE if i % 2 == 0 else NEAR_WHITE
    row_data = [metric, baseline, model, imp]
    for j, val in enumerate(row_data):
        x = Inches(6.8) + j * Inches(1.45)
        add_rect(slide9, x, y, Inches(1.35), Inches(0.5), row_bg)
        color = ACCENT_GREEN if j == 3 and "better" in val else DARK_TEXT
        add_textbox(slide9, x, y + Inches(0.05), Inches(1.35), Inches(0.35),
                    val, font_size=12, bold=(j == 0), color=color, alignment=PP_ALIGN.CENTER)

# Why Ridge
add_textbox(slide9, Inches(6.7), Inches(5.4), Inches(5.8), Inches(0.4),
            "Why Ridge Regression?", font_size=15, bold=True, color=ACCENT_TEAL)
why_ridge = [
    "Interpretable: stakeholders can see which features drive cost",
    "Fast inference: <1ms prediction (critical for real-time demo)",
    "Robust: L2 regularization prevents overfitting on 42 features",
    "Validated: predictions checked against dataset range bounds",
]
for i, reason in enumerate(why_ridge):
    y = Inches(5.85) + i * Inches(0.3)
    add_textbox(slide9, Inches(6.9), y, Inches(5.5), Inches(0.3),
                f"- {reason}", font_size=11, color=DARK_TEXT)

add_speaker_tag(slide9, "Jason")
add_slide_number(slide9, 9)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10: Dynamic Data Analysis  (Jason)
# ══════════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, WHITE)

add_textbox(slide10, Inches(0.6), Inches(0.5), Inches(10), Inches(0.7),
            "Dynamic Data Analysis: Code Generation + Sandbox", font_size=36, bold=True, color=DARK_TEXT)
add_section_header_line(slide10)

# Flow: left to right
da_steps = [
    ("Upload Dataset", "User uploads CSV/Excel\nvia web interface", ACCENT_BLUE),
    ("Metadata\nExtraction", "Extract column names,\ntypes, stats (NOT raw data\n= privacy by design)", ACCENT_TEAL),
    ("Cloud LLM\nCode Gen", "Send metadata to\nZhipu GLM-4 or Gemini\n(dual fallback)", ACCENT_PURPLE),
    ("Docker\nSandbox", "Execute generated Python\nin security-hardened\ncontainer", ACCENT_ORANGE),
    ("Results +\nVisualization", "Return analysis output,\ncharts, and insights\nto user", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(da_steps):
    x = Inches(0.3) + i * Inches(2.55)
    y = Inches(1.8)
    add_rect(slide10, x, y, Inches(2.2), Inches(2.8), color)
    # Step number
    add_textbox(slide10, x + Inches(0.1), y + Inches(0.1), Inches(0.35), Inches(0.35),
                str(i + 1), font_size=18, bold=True, color=WHITE)
    # Title
    for j, line in enumerate(title.split("\n")):
        add_textbox(slide10, x, y + Inches(0.5) + j * Inches(0.3), Inches(2.2), Inches(0.3),
                    line, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Desc
    for j, line in enumerate(desc.split("\n")):
        add_textbox(slide10, x + Inches(0.15), y + Inches(1.3) + j * Inches(0.28), Inches(1.9), Inches(0.3),
                    line, font_size=10, color=RGBColor(0xE8, 0xE8, 0xE8))
    # Arrow
    if i < 4:
        add_textbox(slide10, x + Inches(2.2), y + Inches(1.0), Inches(0.35), Inches(0.4),
                    ">", font_size=20, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom: Security highlights
add_textbox(slide10, Inches(0.6), Inches(5.0), Inches(6), Inches(0.4),
            "Security Hardening (36 TDI Rounds)", font_size=18, bold=True, color=ACCENT_ORANGE)

security_items = [
    "Code validator blocks dangerous imports (os, subprocess, shutil, etc.)",
    "Docker container: no network access, resource limits, read-only filesystem",
    "AST-based static analysis before execution — not just regex pattern matching",
    "Dual cloud LLM fallback: Zhipu -> Gemini -> template-based safe response",
]

for i, item in enumerate(security_items):
    y = Inches(5.5) + i * Inches(0.4)
    # Shield icon
    shield = slide10.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.03),
                                       Inches(0.2), Inches(0.2))
    shield.fill.solid()
    shield.fill.fore_color.rgb = ACCENT_ORANGE
    shield.line.fill.background()
    add_textbox(slide10, Inches(1.15), y, Inches(11), Inches(0.35),
                item, font_size=12, color=DARK_TEXT)

add_speaker_tag(slide10, "Jason")
add_slide_number(slide10, 10)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 11: Results & Evaluation  (Jason)
# ══════════════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide11, WHITE)

add_textbox(slide11, Inches(0.6), Inches(0.5), Inches(10), Inches(0.7),
            "Results & Evaluation", font_size=36, bold=True, color=DARK_TEXT)
add_section_header_line(slide11)

# Three metric cards
result_cards = [
    ("RAG Knowledge QA", [
        "Source citations with document references",
        "Suggested follow-up questions enabled",
        "Hybrid retrieval: BM25 + vector + RRF",
        "Groundedness check on every answer",
        "9+ construction documents indexed",
    ], ACCENT_BLUE, LIGHT_BLUE_BG),

    ("Cost Estimation", [
        "R\u00b2 = 0.989 (vs 0.975 baseline)",
        "MAE: $2.58M (36% improvement)",
        "MAPE: 6.18% (28% improvement)",
        "Reasonableness validation: active",
        "Sub-millisecond inference time",
    ], ACCENT_TEAL, LIGHT_TEAL_BG),

    ("Data Analysis", [
        "Cloud LLM dual fallback: active",
        "Docker sandbox: hardened (36 rounds)",
        "Privacy-by-design: metadata only",
        "Code validator: AST-based",
        "Supports CSV and Excel uploads",
    ], ACCENT_PURPLE, LIGHT_ORANGE_BG),
]

for i, (title, metrics, color, bg) in enumerate(result_cards):
    x = Inches(0.6) + i * Inches(4.2)
    add_rect(slide11, x, Inches(1.8), Inches(3.8), Inches(3.5), bg, border_color=color)
    add_textbox(slide11, x + Inches(0.2), Inches(1.95), Inches(3.4), Inches(0.35),
                title, font_size=16, bold=True, color=color)
    for j, m in enumerate(metrics):
        y = Inches(2.45) + j * Inches(0.52)
        add_textbox(slide11, x + Inches(0.3), y, Inches(3.2), Inches(0.45),
                    f"- {m}", font_size=13, color=DARK_TEXT)

# Bottom: Testing summary
add_rect(slide11, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.3), NEAR_WHITE, border_color=ACCENT_GREEN)
add_textbox(slide11, Inches(0.8), Inches(5.7), Inches(3), Inches(0.35),
            "Testing Summary", font_size=16, bold=True, color=ACCENT_GREEN)

test_stats = [
    ("561", "Unit Tests Collected"),
    ("70%+", "Code Coverage"),
    ("36", "TDI Security Rounds"),
    ("11", "CI Gate Checks"),
]

for i, (num, label) in enumerate(test_stats):
    x = Inches(1.0) + i * Inches(3.0)
    add_textbox(slide11, x, Inches(6.05), Inches(1.5), Inches(0.5),
                num, font_size=28, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
    add_textbox(slide11, x, Inches(6.5), Inches(2.5), Inches(0.3),
                label, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_speaker_tag(slide11, "Jason")
add_slide_number(slide11, 11)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 12: Challenges & Solutions  (Angel)
# ══════════════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide12, WHITE)

add_textbox(slide12, Inches(0.6), Inches(0.5), Inches(10), Inches(0.7),
            "Challenges & How We Solved Them", font_size=36, bold=True, color=DARK_TEXT)
add_section_header_line(slide12)

challenges = [
    ("LLM Performance on Local Hardware",
     "Ollama latency too high for real-time demo",
     "Chose Qwen3.5:4b (28 TPS on M1 Max), disabled thinking mode, added query caching (TTL=120s)",
     ACCENT_BLUE),
    ("Intent Classification Loops",
     "Ambiguous queries caused infinite clarification cycles",
     "Set MAX_CLARIFICATION_ROUNDS=2, confidence threshold=0.8 for direct routing",
     ACCENT_ORANGE),
    ("Code Sandbox Security",
     "User-submitted code could escape sandbox",
     "36 rounds of test-driven security hardening: AST validator + Docker isolation + import blocking",
     ACCENT_PURPLE),
    ("PaddleOCR Compatibility",
     "PaddleOCR requires specific Python version on Apple Silicon",
     "Locked to Python 3.13 + Developer Nightly Build; documented arm64 requirement",
     ACCENT_TEAL),
    ("RAG Answer Quality",
     "Source citations not reliably appearing in early versions",
     "Backend now enforces source fields; frontend rendering improved significantly",
     ACCENT_GREEN),
]

for i, (title, problem, solution, color) in enumerate(challenges):
    y = Inches(1.8) + i * Inches(1.05)
    # Color bar
    bar = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.08), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    # Title
    add_textbox(slide12, Inches(0.9), y, Inches(4.5), Inches(0.3),
                title, font_size=14, bold=True, color=color)
    # Problem
    add_textbox(slide12, Inches(0.9), y + Inches(0.35), Inches(4.5), Inches(0.45),
                f"Problem: {problem}", font_size=11, color=MID_GRAY)
    # Solution
    add_textbox(slide12, Inches(6.5), y + Inches(0.05), Inches(6.2), Inches(0.75),
                f"Solution: {solution}", font_size=12, color=DARK_TEXT)

add_speaker_tag(slide12, "Angel")
add_slide_number(slide12, 12)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 13: Future Work  (Angel)
# ══════════════════════════════════════════════════════════════════════
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide13, WHITE)

add_textbox(slide13, Inches(0.6), Inches(0.5), Inches(10), Inches(0.7),
            "Future Work & Roadmap", font_size=36, bold=True, color=DARK_TEXT)
add_section_header_line(slide13)

future_items = [
    ("Multi-Tenant Production Deployment",
     "Extend X-Tenant-ID isolation to full production with auth, RBAC, and per-tenant data partitioning",
     ACCENT_BLUE),
    ("Cloud-Native Architecture",
     "Containerize all services (Kubernetes), enable horizontal scaling for concurrent users",
     ACCENT_TEAL),
    ("Advanced RAG Techniques",
     "Add graph-based retrieval, multi-modal document understanding, and cross-document reasoning",
     ACCENT_PURPLE),
    ("Expanded ML Models",
     "Add project timeline prediction, resource optimization, safety risk scoring beyond cost estimation",
     ACCENT_ORANGE),
    ("Mobile & Field Access",
     "Responsive UI for tablets/phones so construction teams can query documents on-site",
     ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(future_items):
    y = Inches(1.8) + i * Inches(1.05)
    # Number
    num_bg = slide13.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05),
                                       Inches(0.4), Inches(0.4))
    num_bg.fill.solid()
    num_bg.fill.fore_color.rgb = color
    num_bg.line.fill.background()
    add_textbox(slide13, Inches(0.8), y + Inches(0.07), Inches(0.4), Inches(0.35),
                str(i + 1), font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Title + Desc
    add_textbox(slide13, Inches(1.5), y, Inches(5.5), Inches(0.35),
                title, font_size=16, bold=True, color=color)
    add_textbox(slide13, Inches(1.5), y + Inches(0.4), Inches(10.5), Inches(0.55),
                desc, font_size=13, color=MID_GRAY)

# Right side: vision statement
add_rect(slide13, Inches(8.0), Inches(1.8), Inches(4.8), Inches(3.5), LIGHT_BLUE_BG, border_color=ACCENT_BLUE)
add_textbox(slide13, Inches(8.2), Inches(2.0), Inches(4.4), Inches(0.4),
            "Vision", font_size=20, bold=True, color=ACCENT_BLUE)
add_textbox(slide13, Inches(8.2), Inches(2.5), Inches(4.4), Inches(2.5),
            "Industry AI Flow demonstrates that AI can\n"
            "meaningfully assist the construction industry\n"
            "today — not as a futuristic concept, but as\n"
            "a working, tested prototype that retrieves\n"
            "knowledge, predicts costs, and analyzes\n"
            "data through a single unified platform.",
            font_size=13, color=DARK_TEXT)

add_speaker_tag(slide13, "Angel")
add_slide_number(slide13, 13)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 14: Thank You / Q&A  (Angel)
# ══════════════════════════════════════════════════════════════════════
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide14, WHITE)

# Top accent bar
bar = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

add_textbox(slide14, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
            "Thank You!", font_size=52, bold=True, color=DARK_TEXT,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide14, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.6),
            "Questions & Discussion", font_size=24, color=ACCENT_BLUE,
            alignment=PP_ALIGN.CENTER)

# Divider
div = slide14.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(4.0), Inches(3.3), Inches(0.03)
)
div.fill.solid()
div.fill.fore_color.rgb = LIGHT_GRAY
div.line.fill.background()

# Team
add_multi_text(slide14, Inches(1.5), Inches(4.5), Inches(10.3), Inches(2.0), [
    ("Angel Daniel Bustamante Perez  |  Jason Niu  |  Jack Si", 16, False, DARK_TEXT, PP_ALIGN.CENTER),
    ("", 8),
    ("SAIT  |  Integrated Artificial Intelligence  |  Capstone 2026", 13, False, MID_GRAY, PP_ALIGN.CENTER),
    ("", 8),
    ("Instructor: Reeta", 13, False, MID_GRAY, PP_ALIGN.CENTER),
])

# References section
add_textbox(slide14, Inches(0.6), Inches(5.8), Inches(12), Inches(0.35),
            "References", font_size=14, bold=True, color=ACCENT_BLUE)

refs = [
    "McKinsey & Company (2017). Reinventing Construction: A Route to Higher Productivity.",
    "LangChain / LangGraph — https://www.langchain.com/langgraph",
    "Ollama — https://ollama.com | Qwen 3.5 — https://huggingface.co/Qwen",
    "pgvector — https://github.com/pgvector/pgvector | BGE Reranker — https://huggingface.co/BAAI/bge-reranker-base",
    "PaddleOCR — https://github.com/PaddlePaddle/PaddleOCR | FastAPI — https://fastapi.tiangolo.com",
]
for i, ref in enumerate(refs):
    add_textbox(slide14, Inches(0.8), Inches(6.15) + i * Inches(0.22), Inches(11.5), Inches(0.25),
                ref, font_size=9, color=MID_GRAY)

add_slide_number(slide14, 14)


# ── Save ─────────────────────────────────────────────────────────────
output_path = "docs/presentation/Industry_AI_Flow_Capstone_Presentation.pptx"
import os
os.makedirs("docs/presentation", exist_ok=True)
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
